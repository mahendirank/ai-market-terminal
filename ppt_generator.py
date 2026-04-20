from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

OUTPUT_DIR = os.path.expanduser("~/ai-system/outputs")

def create_ppt(text):
    prs = Presentation()

    slides = text.split("SLIDE:")

    for slide_data in slides:
        slide_data = slide_data.strip()
        if not slide_data:
            continue

        lines = slide_data.split("\n")
        title = lines[0].strip()

        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Title styling
        slide.shapes.title.text = title
        title_shape = slide.shapes.title.text_frame.paragraphs[0]
        title_shape.font.size = Pt(32)
        title_shape.font.bold = True

        # Content
        content = slide.placeholders[1].text_frame
        content.clear()

        for line in lines[1:]:
            if line.strip():
                p = content.add_paragraph()
                p.text = line.replace("-", "").strip()
                p.level = 0
                p.font.size = Pt(18)

    file_path = os.path.join(OUTPUT_DIR, "presentation_pro.pptx")
    prs.save(file_path)

    print(f"\n✅ PRO PPT Created: {file_path}")
